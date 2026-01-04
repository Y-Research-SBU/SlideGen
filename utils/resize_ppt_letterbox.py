#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Resize PPTX to a new aspect ratio using 'Ensure Fit' behavior (no stretch):
1) Neutralize per-slide background to avoid 16:9 master visuals breaking layout after resize.
2) Uniformly scale and center all shapes (contain), preserving aspect ratio (no distortion).

Usage:
  pip install python-pptx
  python resize_ppt_ensure_fit.py input.pptx output_4x3.pptx --width 10 --ratio 4:3 --bg 255,255,255

Arguments:
  --width : Target slide width in inches (height is computed from ratio)
  --ratio : Target aspect ratio "W:H" (e.g., 4:3, 16:10, 5:4, ...)
  --bg    : Background color for each slide in R,G,B (default 255,255,255)

Notes:
- This script does NOT edit the Slide Master (python-pptx limitation).
  It neutralizes the slide background (solid fill) so old 16:9 master art doesn't conflict.
- If you need a texture/graphic background, convert that into a per-slide bottom picture
  and keep 'Ensure Fit' for foreground content to avoid distortion.


export SOFFICE_PATH="SlideGen/LibreOffice_25.2.4.3_Linux_x86-64_deb/libreoffice-extracted/opt/libreoffice25.2/program/soffice"

PATH="SlideGen/LibreOffice_25.2.4.3_Linux_x86-64_deb/libreoffice-extracted/opt/libreoffice25.2/program:$PATH" \

python SlideGen/utils/resize_ppt_letterbox.py \
    --input SlideGen/contents/AlphaFold_Meets_Flow_Matching_for_Generating_Protein_Ensembles/4o_4o_output_slides.pptx\
    --output SlideGen/contents/AlphaFold_Meets_Flow_Matching_for_Generating_Protein_Ensembles/4o_4o_output_slides43.pptx

/opt/conda/envs/paper2pptx/bin/python SlideGen/utils/resize_ppt_letterbox.py \
  --input  SlideGen/contents/AlphaFold_Meets_Flow_Matching_for_Generating_Protein_Ensembles/4o_4o_output_slides.pptx \
  --output SlideGen/contents/AlphaFold_Meets_Flow_Matching_for_Generating_Protein_Ensembles/4o_4o_output_slides43.pptx

  
"""

import argparse
import os
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor


# -----------------------------
# Utilities
# -----------------------------
def _run(cmd):
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"Command failed: {' '.join(cmd)}\nSTDOUT:\n{proc.stdout}\nSTDERR:\n{proc.stderr}")
    return proc


def _which(bin_name):
    return shutil.which(bin_name) is not None


def _parse_ratio(ratio_str: str):
    try:
        rw, rh = ratio_str.split(":")
        return float(rw), float(rh)
    except Exception:
        raise ValueError('Invalid --ratio. Use format like "4:3" or "16:10".')


def _parse_bg(bg_str: str):
    try:
        parts = [int(x) for x in bg_str.split(",")]
        if len(parts) != 3 or any(not (0 <= v <= 255) for v in parts):
            raise ValueError
        return tuple(parts)
    except Exception:
        raise ValueError('Invalid --bg. Use "R,G,B" with values 0-255, e.g., "245,245,245".')


# -----------------------------
# Safe mode (rasterize): LO -> PDF -> PNG -> paste
# -----------------------------
def convert_safe_raster(input_path, output_path, width_in=10.0, ratio_w=4.0, ratio_h=3.0, bg_rgb=(255, 255, 255), dpi=220):
    """
    Visually lossless 'Ensure Fit':
    1) Convert PPTX to PDF via LibreOffice headless
    2) Convert PDF pages to PNG via Poppler (pdftoppm)
    3) Create a new PPTX and, per slide, place the PNG centered & letterboxed at target size
    """
    # Dependencies
    lo_bin = "soffice"  # libreoffice binary
    have_lo = _which(lo_bin)
    have_pdftoppm = _which("pdftoppm")
    if not have_lo:
        raise RuntimeError("LibreOffice (soffice) not found. Install it first.")
    if not have_pdftoppm:
        raise RuntimeError("Poppler (pdftoppm) not found. Install it first.")

    with tempfile.TemporaryDirectory() as tmpd:
        tmp_dir = Path(tmpd)
        # 1) PPTX -> PDF
        pdf_out = tmp_dir / "deck.pdf"
        _run([lo_bin, "--headless", "--convert-to", "pdf", "--outdir", str(tmp_dir), str(input_path)])
        # LibreOffice names the output by replacing extension with .pdf
        # Ensure expected filename
        guessed_pdf = tmp_dir / (Path(input_path).stem + ".pdf")
        if guessed_pdf.exists():
            guessed_pdf.rename(pdf_out)
        if not pdf_out.exists():
            raise RuntimeError("Failed to produce PDF via LibreOffice.")

        # 2) PDF -> PNG pages
        # pdftoppm -png -r <dpi> deck.pdf slide
        _run(["pdftoppm", "-png", "-r", str(dpi), str(pdf_out), str(tmp_dir / "slide")])
        # Now we have slide-1.png, slide-2.png, ...
        pngs = sorted(tmp_dir.glob("slide-*.png"), key=lambda p: int(p.stem.split("-")[-1]))
        if not pngs:
            raise RuntimeError("No PNG pages generated from PDF.")

        # 3) Build new deck and place images
        prs = Presentation()
        # Set target size
        new_w = Inches(width_in)
        new_h = Emu(new_w * ratio_h / ratio_w)
        prs.slide_width, prs.slide_height = new_w, new_h

        # Blank layout
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        # Background color neutralization
        r, g, b = bg_rgb

        for png in pngs:
            slide = prs.slides.add_slide(blank_layout)
            # Solid background
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(r, g, b)

            # Add picture at 0,0 then compute scale-to-fit
            pic = slide.shapes.add_picture(str(png), Emu(0), Emu(0))
            img_w, img_h = pic.width, pic.height

            sx = float(new_w) / float(img_w)
            sy = float(new_h) / float(img_h)
            s = min(sx, sy)  # contain
            w2 = Emu(img_w * s)
            h2 = Emu(img_h * s)
            dx = Emu((float(new_w - w2) / 2.0))
            dy = Emu((float(new_h - h2) / 2.0))

            # Reposition & resize
            pic.left = dx
            pic.top = dy
            pic.width = w2
            pic.height = h2

        prs.save(output_path)


# -----------------------------
# Optional: best-effort vector mode (may still break on masters/placeholders)
# -----------------------------
def _scale_and_center_shape(shape, s, dx, dy):
    shape.left = Emu(shape.left * s + dx)
    shape.top = Emu(shape.top * s + dy)
    shape.width = Emu(shape.width * s)
    shape.height = Emu(shape.height * s)

    # Try to curb auto-fit surprises for text frames (helps a bit, not bulletproof)
    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        tf = shape.text_frame
        try:
            tf.word_wrap = True
        except Exception:
            pass
        # Do not force font size; PowerPoint can still reflow text based on new shape box.


def _process_group_shape(group, s, dx, dy):
    _scale_and_center_shape(group, s, dx, dy)
    for shp in group.shapes:
        shp.left = Emu(shp.left * s)
        shp.top = Emu(shp.top * s)
        shp.width = Emu(shp.width * s)
        shp.height = Emu(shp.height * s)


def _neutralize_background(prs, bg_rgb=(255, 255, 255)):
    r, g, b = bg_rgb
    for slide in prs.slides:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)


def convert_best_effort_vector(input_path, output_path, width_in=10.0, ratio_w=4.0, ratio_h=3.0, bg_rgb=(255, 255, 255)):
    prs = Presentation(input_path)
    old_w, old_h = prs.slide_width, prs.slide_height

    _neutralize_background(prs, bg_rgb=bg_rgb)

    new_w = Inches(width_in)
    new_h = Emu(new_w * ratio_h / ratio_w)
    prs.slide_width, prs.slide_height = new_w, new_h

    sx = float(new_w) / float(old_w)
    sy = float(new_h) / float(old_h)
    s = min(sx, sy)
    dx = float(new_w - old_w * s) / 2.0
    dy = float(new_h - old_h * s) / 2.0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _process_group_shape(shape, s, dx, dy)
            else:
                _scale_and_center_shape(shape, s, dx, dy)

    prs.save(output_path)


# -----------------------------
# CLI
# -----------------------------
if __name__ == "__main__":
    ap = argparse.ArgumentParser(description="Resize PPTX with reliable 'Ensure Fit'. Default uses raster-safe pipeline.")
    ap.add_argument("--input", required=True, help="Input .pptx")
    ap.add_argument("--output", required=True, help="Output .pptx")
    ap.add_argument("--width", type=float, default=10.0, help="Target slide width in inches (default: 10.0)")
    ap.add_argument("--ratio", type=str, default="4:3", help='Target aspect ratio "W:H" (default: 4:3)')
    ap.add_argument("--bg", type=str, default="255,255,255", help='Background color "R,G,B" (default: 255,255,255)')
    ap.add_argument("--dpi", type=int, default=220, help="Rasterization DPI (default: 220)")
    ap.add_argument(
        "--mode",
        choices=["safe", "vector"],
        default="safe",
        help="safe: render to image then paste (recommended); vector: best-effort shape scaling",
    )
    args = ap.parse_args()

    ratio_w, ratio_h = _parse_ratio(args.ratio)
    bg_rgb = _parse_bg(args.bg)

    if args.mode == "safe":
        convert_safe_raster(
            args.input, args.output,
            width_in=args.width,
            ratio_w=ratio_w,
            ratio_h=ratio_h,
            bg_rgb=bg_rgb,
            dpi=args.dpi,
        )
    else:
        convert_best_effort_vector(
            args.input, args.output,
            width_in=args.width,
            ratio_w=ratio_w,
            ratio_h=ratio_h,
            bg_rgb=bg_rgb,
        )
