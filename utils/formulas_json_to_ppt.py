# -*- coding: utf-8 -*-
"""
Render LaTeX formulas from a JSON file into PNG images and insert them into a PPT.

Input JSON format example (array of objects):
[
  {
    "page": 2,
    "bbox": {"x": 123.4, "y": 567.8, "w": 90.1, "h": 22.3},
    "latex_raw": "\\mathcal { L } _ { d i s t i l l } = ...",
    "latex": "\\mathcal{L}_{distill} = ...",
    "crop_path": ""
  },
  ...
]

Dependencies:
    pip install matplotlib python-pptx pillow

If you have a TeX distribution installed (TeX Live / MiKTeX), you can pass --usetex
for higher-quality rendering.
"""

import json,re
import argparse
from pathlib import Path
import matplotlib
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

import matplotlib
import matplotlib.pyplot as plt
from pathlib import Path
import shutil
HAS_TEX = shutil.which("latex") is not None  # 装了 TeX 才能用 usetex

def _render_once(latex: str, out_path: Path, dpi: int, fontsize: int, usetex: bool):
    
    try:
        matplotlib.rcParams["text.parse_math"] = True
    except Exception:
        pass
    matplotlib.rcParams.update({
        "text.usetex": (usetex and HAS_TEX),   
        "mathtext.fontset": "cm",              
        "font.size": fontsize,
        "figure.dpi": dpi,
        "figure.max_open_warning": 0,          
    })
     
    fig = plt.figure(figsize=(0.01, 0.01))
    try:
        ax = fig.add_subplot(111)
        ax.axis("off")
        txt = latex if (latex.startswith("$") and latex.endswith("$")) else f"${latex}$"
        ax.text(0.5, 0.5, txt, ha="center", va="center")
        out_path.parent.mkdir(parents=True, exist_ok=True)
        fig.savefig(out_path, dpi=dpi, bbox_inches="tight", pad_inches=0.12, transparent=True)
    finally:
        plt.close(fig)   

 
_SP_FIX = re.compile(
    r"(?:(?<=\\)\s+)|(?<=_)\s+|(?<=\^)\s+|(?<=\{)\s+|\s+(?=\})|\s+(?=[,=)])|(?<=\()\s+"
)

# \cmd { → \cmd{
_CMD_BRACE_FIX = re.compile(r"\\([A-Za-z]+)\s+\{")

def _collapse_spaces_inside_braces(text: str) -> str:
    # 把 { v i t } → {vit}， 
    def repl(m):
        inner = m.group(1)
        inner2 = re.sub(r"(?<=\w)\s+(?=\w)", "", inner)
        return "{" + inner2 + "}"
    return re.sub(r"\{([^{}]+)\}", repl, text)

def normalize_formula_strong(s: str) -> str:
    s = (s or "").strip()
    if not s:
        return s
 
    s = _SP_FIX.sub("", s)
    s = _CMD_BRACE_FIX.sub(r"\\\1{", s)
    s = _collapse_spaces_inside_braces(s)
 
    s = re.sub(r"h_\{\\amalgm\}", r"h_{llm}", s)
    s = re.sub(r"h_\{\\amalg\s*m\}", r"h_{llm}", s)
    
    s = re.sub(r"\\amalgm", "llm", s)
    s = re.sub(r"\\amalg\s*m", "llm", s)
    s = re.sub(r"\\amalg\b", "ll", s)
  
    # Case A: \mathrm\vec{vit} -> \vec{vit}  (keep the vector notation)
    s = re.sub(r"\\mathrm\\vec\{([^{}]+)\}", r"\\vec{\1}", s)

    # Case B: \mathrm{\vec{vit}} -> \vec{vit}
    s = re.sub(r"\\mathrm\{\\vec\{([^{}]+)\}\}", r"\\vec{\1}", s)

    # Case C: \mathrmvit -> \mathrm{vit}
    s = re.sub(r"\\mathrm([A-Za-z]+)", r"\\mathrm{\1}", s)

    # Case D: \vecvit -> \vec{vit}
    s = re.sub(r"\\vec([A-Za-z]+)", r"\\vec{\1}", s)


    # 3) \logP → \log P
    s = re.sub(r"\\log([A-Za-z])", r"\\log \1", s)

    # 4) 多个 \, 合并为一个
    s = re.sub(r"(\\,){2,}", r"\\,", s)

    s = re.sub(r"\\mathbf([A-Za-z])", r"\\mathbf{\1}", s)
    s = re.sub(r"\\mathbf([A-Za-z])", r"\\mathbf{\1}", s)
    s = re.sub(r"\\mathrm([A-Za-z])", r"\\mathrm{\1}", s)
    s = re.sub(r"\\mathrm([A-Za-z])", r"\\mathrm{\1}", s)   
 
    s = re.sub(r"\\mathbb\{([A-Za-z])\}", r"\\mathrm{\1}", s)
 
    s = s.replace(r"\logp",  r"\log p")
    s = s.replace(r"\circp", r"\circ p")
    s = s.replace(r"\cdotf", r"\cdot f")
    s = s.replace(r"\coloneqq", ":=")
 
    s = re.sub(r"\\for\b", r"\\mathrm{for}", s)
 
    s = re.sub(r"\\mathrm\\vec\{([^{}]+)\}", r"\\vec{\1}", s)
    s = re.sub(r"\\mathrm\{\\vec\{([^{}]+)\}\}", r"\\vec{\1}", s)
 
    s = re.sub(r"\\simD\b", r"\\sim\\mathcal{D}", s)
 
    s = re.sub(r"\\mathbf([A-Za-z])", r"\\mathbf{\1}", s)
    s = re.sub(r"\\mathrm([A-Za-z])", r"\\mathrm{\1}", s)
 
    s = re.sub(r"\\mathbb\{([A-Za-z])\}", r"\\mathrm{\1}", s)
 
    s = s.replace(r"\logp",  r"\log p")
    s = s.replace(r"\circp", r"\circ p")
    s = s.replace(r"\cdotf", r"\cdot f")
    s = s.replace(r"\coloneqq", ":=")
    s = re.sub(r"\\for\b", r"\\mathrm{for}", s)   # \for → ‘for’
 
    s = re.sub(r"h_\{\\amalgm\}", r"h_{llm}", s)
    s = re.sub(r"\\amalgm",       "llm",      s)
    s = re.sub(r"\\amalg\s*m",    "llm",      s)
    s = re.sub(r"\\amalg\b",      "ll",       s)
 
    s = re.sub(r"\\mathrm\\vec\{([^{}]+)\}",   r"\\vec{\1}", s)
    s = re.sub(r"\\mathrm\{\\vec\{([^{}]+)\}\}", r"\\vec{\1}", s)
 
    s = re.sub(r"\\nabla_\{\\mathrm\{([A-Za-z])\}\}_\{([^{}]+)\}", r"\\nabla_{\1_{\2}}", s)
    s = re.sub(r"\\nabla_\{\\mathrm\{([A-Za-z])\}\}_([A-Za-z0-9]+)", r"\\nabla_{\1_{\2}}", s)
    s = re.sub(r"\\nabla_\{([A-Za-z])\}_\{([^{}]+)\}", r"\\nabla_{\1_{\2}}", s)  # 更通用
 
    s = re.sub(r"\\simD\b", r"\\sim\\mathcal{D}", s)
    s = re.sub(r"\\simq\b", r"\\sim q", s)
 
    s = re.sub(r"\\underbrace\{([^{}]+)\}\_\{[^{}]*\}", r"\1", s)
 
    s = re.sub(r"[_\\]$", "", s)
 
    n_left  = len(re.findall(r"\\left\(", s))
    n_right = len(re.findall(r"\\right\)", s))
    if n_left > n_right:
        s += "\\right)" * (n_left - n_right)
 
    diff_paren = s.count("(") - s.count(")")
    if diff_paren > 0:
        s += ")" * diff_paren
    diff_brace = s.count("{") - s.count("}")
    if diff_brace > 0:
        s += "}" * diff_brace
 
    s = s.replace(r"\simD", r"\sim\mathcal{D}")
    s = s.replace(r"\simq", r"\sim q")
 
    s = re.sub(r"\\underbrace\s*\{([^{}]+)\}\s*_\s*\{[^{}]*\}", r"\1", s)
    s = re.sub(r"\\underbrace\s*\{([^{}]+)\}(\s*_\s*\{[^{}]*\})?", r"\1", s)
 
    s = re.sub(r"\\nabla_\{\\mathrm\{([A-Za-z])\}\}_\{([^{}]+)\}", r"\\nabla_{\1_{\2}}", s)
    s = re.sub(r"\\nabla_\{([A-Za-z])\}_\{([^{}]+)\}", r"\\nabla_{\1_{\2}}", s)
    s = re.sub(r"\\nabla_\{([A-Za-z])\}_([A-Za-z0-9]+)", r"\\nabla_{\1_{\2}}", s)
 
    s = re.sub(r"\\mathbf([A-Za-z])", r"\\mathbf{\1}", s)
 
    s = re.sub(r"([_^]|\\)\s*$", "", s)
 
    def _balance(pair_left, pair_right):
        nL = len(re.findall(pair_left, s))
        nR = len(re.findall(pair_right, s))
        return nL - nR
    diff = _balance(r"\\left\(", r"\\right\)")
    if diff > 0: s += "\\right)" * diff
    diff = _balance(r"\\left\[", r"\\right\]")
    if diff > 0: s += "\\right]" * diff
    diff = _balance(r"\\left\\\{", r"\\right\\\}")
    if diff > 0: s += "\\right\\}" * diff
    diff = _balance(r"\\left\\\|", r"\\right\\\|")
    if diff > 0: s += "\\right\\|" * diff
 
    diff_paren = s.count("(") - s.count(")")
    if diff_paren > 0: s += ")" * diff_paren
    diff_brace = s.count("{") - s.count("}")
    if diff_brace > 0: s += "}" * diff_brace
 
    s = re.sub(r"\s+", " ", s).strip()
   
    return s
def render_latex_to_png(latex: str, out_path: Path, dpi: int = 220, fontsize: int = 24, usetex: bool = False):
     
    latex_norm = normalize_formula_strong(latex)
    try:
        _render_once(latex_norm, out_path, dpi, fontsize, usetex=False)
        return
    except Exception as e1:
        if usetex:
           
            _render_once(latex_norm, out_path, dpi, fontsize, usetex=True)
            return
        else:
            
            _render_once(latex_norm, out_path, dpi, fontsize, usetex=True)
            return



def formulas_json_to_ppt(
    json_path: Path,
    ppt_path: Path,
    img_dir: Path,
    dpi: int = 220,
    fontsize: int = 24,
    usetex: bool = False,
    max_per_slide: int = 1,
    margin_in: float = 1.0,
):
    """
    Read formulas from JSON; render each as PNG; insert into a PPT.

    Args:
        json_path: path to JSON file
        ppt_path: output PPTX path
        img_dir: directory to store rendered images
        dpi: PNG dpi
        fontsize: LaTeX font size (points)
        usetex: use system LaTeX (needs TeX Live/MiKTeX); fallback is mathtext
        max_per_slide: how many formulas per slide (1 = one per slide)
        margin_in: margin (inches) from slide borders
    """
    data = json.loads(Path(json_path).read_text(encoding="utf-8"))

    # Accept either a list or an object with "items"
    if isinstance(data, dict) and "items" in data:
        items = data["items"]
    else:
        items = data

    # Build list of latex strings
    formulas = []
    for it in items:
        latex = it.get("latex") or it.get("latex_raw") or ""
        latex = latex.strip()
        if latex:
            formulas.append(latex)

    if not formulas:
        print("No formulas found in JSON (looked for 'latex' or 'latex_raw').")
        return

    prs = Presentation()
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Convert EMU to inches helpers (pptx stores sizes in EMUs)
    EMU_PER_INCH = 914400
    slide_w_in = slide_w / EMU_PER_INCH
    slide_h_in = slide_h / EMU_PER_INCH

    img_dir.mkdir(parents=True, exist_ok=True)

    # Render and place
    slide = None
    per_slide_count = 0

    for idx, latex in enumerate(formulas, start=1):
        # 1) Render PNG
        img_path = img_dir / f"formula_{idx:04d}.png"
        try:
            print("input :", latex)
            render_latex_to_png(latex, img_path, dpi=dpi, fontsize=fontsize, usetex=usetex)
        except Exception as e:
            print(f"[Warning] Failed to render formula #{idx}: {e}")
            continue

        # 2) Add slide when needed
        if slide is None or per_slide_count >= max_per_slide:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            per_slide_count = 0

        # 3) Compute placement
        # Try to keep image within slide, centered; width = slide_w - 2*margin
        content_w_in = max(0.5, slide_w_in - 2 * margin_in)
        left_in = margin_in
        # default top: center vertically if one per slide; otherwise stack
        if max_per_slide == 1:
            top_in = (slide_h_in - content_w_in * 0.5) / 2  # rough center guess
            if top_in < margin_in:
                top_in = margin_in
        else:
            # stack with equal spacing
            row_h_in = (slide_h_in - 2 * margin_in) / max_per_slide
            top_in = margin_in + per_slide_count * row_h_in + (row_h_in - row_h_in * 0.6) / 2
            content_w_in = content_w_in  # keep width

        # Insert picture; pptx will scale height to preserve aspect ratio if only width given
        pic = slide.shapes.add_picture(
            str(img_path),
            left=int(left_in * EMU_PER_INCH),
            top=int(top_in * EMU_PER_INCH),
            width=int(content_w_in * EMU_PER_INCH),
            height=None,
        )

        # Optional: add a small caption text placeholder (commented)
        # tx = slide.shapes.add_textbox(
        #     left=int(margin_in * EMU_PER_INCH),
        #     top=int((top_in + (pic.height/EMU_PER_INCH) + 0.2) * EMU_PER_INCH),
        #     width=int((slide_w_in - 2*margin_in) * EMU_PER_INCH),
        #     height=int(0.5 * EMU_PER_INCH),
        # )
        # tx.text_frame.text = f"Formula {idx}"

        per_slide_count += 1

    prs.save(ppt_path)
    print(f"Done. Rendered {len(formulas)} formulas.")
    print(f"PPT saved to: {ppt_path}")
    print(f"PNGs saved under: {img_dir}")

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--json", required=True, help="Path to formulas JSON (array of items)")
    ap.add_argument("--out-ppt", default="formulas.pptx", help="Output PPTX path")
    ap.add_argument("--img-dir", default="formula_imgs", help="Directory to save rendered PNGs")
    ap.add_argument("--dpi", type=int, default=220, help="PNG DPI (default 220)")
    ap.add_argument("--fontsize", type=int, default=24, help="LaTeX font size (pt)")
    ap.add_argument("--usetex", action="store_true", help="Use system LaTeX (needs TeX Live/MiKTeX)")
    ap.add_argument("--max-per-slide", type=int, default=1, help="Formulas per slide (default 1)")
    ap.add_argument("--margin-in", type=float, default=1.0, help="Slide margin in inches (default 1.0)")
    args = ap.parse_args()

    formulas_json_to_ppt(
        json_path=Path(args.json),
        ppt_path=Path(args.out_ppt),
        img_dir=Path(args.img_dir),
        dpi=args.dpi,
        fontsize=args.fontsize,
        usetex=args.usetex,
        max_per_slide=args.max_per_slide,
        margin_in=args.margin_in,
    )

if __name__ == "__main__":
    main()



'''
 
python formulas_json_to_ppt.py \
  --json formulas_with_bbox.json \
  --out-ppt formulas.pptx \
  --img-dir formula_imgs \
  --dpi 220 \
  --fontsize 26

 
python formulas_json_to_ppt.py \
  --json formulas_with_bbox.json \
  --out-ppt formulas_usetex.pptx \
  --img-dir formula_imgs_tex \
  --dpi 250 \
  --fontsize 26 \
  --usetex
 

'''