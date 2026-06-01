"""Template introspection for SlideGen v2.

Lets users drop an arbitrary ``.pptx`` into ``utils/slides_template/`` and have
the pipeline use it, instead of relying on the fixed 24-layout master deck with
hard-coded layout/placeholder names.

The public surface is:

* :func:`scan_template`            -> :class:`TemplateSpec`
* :func:`build_layout_library_md`  -> markdown table injected into the arranger prompt
* :func:`resolve_template_path`    -> locate the chosen ``.pptx``

Role inference is heuristic (placeholder type + position + name hints). A
``<stem>.slidegen.json`` sidecar next to the template can override anything the
heuristics get wrong, so no naming convention is forced on the user.
"""

from __future__ import annotations

import json
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE


# Placeholder types that carry text we fill (title / part-number / bullets).
# DATE / FOOTER / SLIDE_NUMBER are template chrome and explicitly excluded.
TEXT_PH_TYPES = {
    PH_TYPE.TITLE,
    PH_TYPE.CENTER_TITLE,
    PH_TYPE.SUBTITLE,
    PH_TYPE.BODY,
    PH_TYPE.OBJECT,  # some decks use generic object placeholders for body text
}
CHROME_PH_TYPES = {PH_TYPE.DATE, PH_TYPE.FOOTER, PH_TYPE.SLIDE_NUMBER}

_PIC_NAME_HINTS = ("picture", "图片", "圖片")

DEFAULT_TEMPLATE_DIR = "utils/slides_template"


# --------------------------------------------------------------------------- #
# Data model
# --------------------------------------------------------------------------- #
@dataclass
class PlaceholderInfo:
    idx: int
    name: str
    type: str           # human-readable placeholder type name
    left: int           # EMU
    top: int
    width: int
    height: int

    @property
    def area(self) -> int:
        return max(0, self.width) * max(0, self.height)


@dataclass
class LayoutSpec:
    name: str
    text_phs: List[PlaceholderInfo] = field(default_factory=list)
    pic_phs: List[PlaceholderInfo] = field(default_factory=list)

    # Inferred content roles (placeholder idx values). None when absent.
    part_idx: Optional[int] = None       # tiny section-number box
    title_idx: Optional[int] = None      # the subsection/section title bar
    body_idxs: List[int] = field(default_factory=list)   # bullet boxes, reading order
    pic_idxs: List[int] = field(default_factory=list)    # picture slots, reading order
    # Column model: each column is {"title": idx|None, "body": idx}. For simple
    # single-body layouts this is one column with title=None; multi-column
    # layouts (e.g. T19) pair a short heading box with its tall bullet box.
    body_cols: List[dict] = field(default_factory=list)

    @property
    def n_body(self) -> int:
        return len(self.body_cols) if self.body_cols else len(self.body_idxs)

    @property
    def n_pic(self) -> int:
        return len(self.pic_idxs)

    def ph_by_idx(self, idx: Optional[int]) -> Optional[PlaceholderInfo]:
        if idx is None:
            return None
        for ph in (*self.text_phs, *self.pic_phs):
            if ph.idx == idx:
                return ph
        return None

    def title_type_idx(self) -> Optional[int]:
        """idx of the most title-like text placeholder (CENTER_TITLE/TITLE first)."""
        for want in ("CENTER_TITLE", "TITLE"):
            for ph in self.text_phs:
                if ph.type == want:
                    return ph.idx
        if self.title_idx is not None:
            return self.title_idx
        if self.text_phs:
            return sorted(self.text_phs, key=lambda p: (p.top, p.left))[0].idx
        return None

    def secondary_text_idxs(self, exclude: Optional[int] = None) -> List[int]:
        """idx list of non-title text placeholders (subtitle/body), reading order."""
        rest = [p for p in self.text_phs if p.idx != exclude]
        rest.sort(key=lambda p: (p.top, p.left))
        return [p.idx for p in rest]


@dataclass
class TemplateSpec:
    path: str
    slide_width: int          # EMU
    slide_height: int
    layouts: Dict[str, LayoutSpec] = field(default_factory=dict)

    # role -> layout name for the four special slides
    roles: Dict[str, Optional[str]] = field(default_factory=dict)
    # content layout names selectable by the arranger
    content_layouts: List[str] = field(default_factory=list)

    def get(self, layout_name: str) -> Optional[LayoutSpec]:
        return self.layouts.get(layout_name)

    def role_layout(self, role: str) -> Optional[str]:
        return self.roles.get(role)


# --------------------------------------------------------------------------- #
# Scanning
# --------------------------------------------------------------------------- #
def _ph_type_name(shape) -> str:
    try:
        return PH_TYPE(shape.placeholder_format.type).name
    except (ValueError, TypeError):
        return str(getattr(shape.placeholder_format, "type", "UNKNOWN"))


def _is_picture_ph(shape) -> bool:
    try:
        if shape.placeholder_format.type == PH_TYPE.PICTURE:
            return True
    except (ValueError, TypeError):
        pass
    name = (getattr(shape, "name", "") or "").lower()
    return any(h in name for h in _PIC_NAME_HINTS)


def _collect_placeholders(layout) -> tuple[List[PlaceholderInfo], List[PlaceholderInfo]]:
    text_phs: List[PlaceholderInfo] = []
    pic_phs: List[PlaceholderInfo] = []
    for sh in layout.placeholders:
        pf = sh.placeholder_format
        try:
            ph_type = pf.type
        except (ValueError, TypeError):
            ph_type = None

        info = PlaceholderInfo(
            idx=pf.idx,
            name=sh.name or "",
            type=_ph_type_name(sh),
            left=sh.left or 0,
            top=sh.top or 0,
            width=sh.width or 0,
            height=sh.height or 0,
        )

        if _is_picture_ph(sh):
            pic_phs.append(info)
        elif ph_type in CHROME_PH_TYPES:
            continue  # skip date/footer/slide-number chrome
        elif ph_type in TEXT_PH_TYPES:
            text_phs.append(info)
        # anything else (e.g. CHART/TABLE/MEDIA) is ignored for now
    return text_phs, pic_phs


def _assign_content_roles(layout: LayoutSpec, slide_width: int) -> None:
    """Infer part / title / body roles for a content layout by geometry.

    Mirrors the original ``find_text_placeholders`` heuristic: the two
    top-most text boxes form a row (left=part number, right=title bar), and
    everything below is body text. Works for single- and multi-body layouts.
    """
    phs = sorted(layout.text_phs, key=lambda p: (p.top, p.left))
    if not phs:
        return

    # A "title row" is the (up to) two shallowest boxes that are short in
    # height relative to the slide (title bars / part numbers are thin).
    title_row_h = 0.18 * (layout_height_guess(layout) or 1)

    if len(phs) >= 3:
        top_two = sorted(phs[:2], key=lambda p: p.left)
        # part = the small box (narrow width), title = the wide one
        a, b = top_two
        if a.width <= b.width:
            layout.part_idx, layout.title_idx = a.idx, b.idx
        else:
            layout.part_idx, layout.title_idx = b.idx, a.idx
        body = phs[2:]
    elif len(phs) == 2:
        # title + single body
        layout.title_idx = phs[0].idx
        body = phs[1:]
    else:
        # single text box -> treat as body
        body = phs

    # Order body boxes in reading order: by column (left) then row (top).
    body_sorted = sorted(body, key=lambda p: (round(p.left / max(slide_width, 1) * 3), p.top))
    layout.body_idxs = [p.idx for p in body_sorted]
    layout.body_cols = _group_body_columns(body, slide_width)

    layout.pic_idxs = [p.idx for p in sorted(layout.pic_phs, key=lambda p: (p.top, p.left))]


def _group_body_columns(body: List[PlaceholderInfo], slide_width: int) -> List[dict]:
    """Group body placeholders into columns of {"title": idx|None, "body": idx}.

    Boxes are bucketed by horizontal band (left position). Within a column, a
    short box stacked above a tall box is treated as (heading, bullets); a lone
    box is the bullets with no heading. This recovers the two-column "heading +
    text" structure of layouts like T19 without hard-coding placeholder ids.
    """
    if not body:
        return []
    # Bucket by left into columns (3 bands across the slide width).
    buckets: Dict[int, List[PlaceholderInfo]] = {}
    for p in body:
        band = round(p.left / max(slide_width, 1) * 3)
        buckets.setdefault(band, []).append(p)

    cols: List[dict] = []
    for band in sorted(buckets):
        boxes = sorted(buckets[band], key=lambda p: p.top)
        if len(boxes) == 1:
            cols.append({"title": None, "body": boxes[0].idx})
            continue
        # Heading = the short top box; bullets = the tallest box below it.
        top = boxes[0]
        tall = max(boxes, key=lambda p: p.height)
        if tall.idx == top.idx:
            # top box is also the tallest -> no separate heading
            cols.append({"title": None, "body": top.idx})
            for extra in boxes[1:]:
                cols.append({"title": None, "body": extra.idx})
        else:
            cols.append({"title": top.idx, "body": tall.idx})
            for extra in boxes:
                if extra.idx not in (top.idx, tall.idx):
                    cols.append({"title": None, "body": extra.idx})
    return cols


def layout_height_guess(layout: LayoutSpec) -> int:
    bottoms = [p.top + p.height for p in (*layout.text_phs, *layout.pic_phs)]
    return max(bottoms) if bottoms else 0


# --------------------------------------------------------------------------- #
# Role inference (cover / toc / divider / last / content)
# --------------------------------------------------------------------------- #
# Name hints recognized for the four special layouts (case-insensitive
# substrings). These are *hints* — geometry decides when names are unknown.
_ROLE_NAME_HINTS = {
    "cover":   ("title slide", "cover", "封面", "标题幻灯片"),
    "toc":     ("mulu", "目录", "content", "agenda", "outline", "toc"),
    "divider": ("dan_mulu", "danmulu", "section", "divider", "章节", "分节"),
    "last":    ("last", "thanks", "thank you", "end", "结束", "致谢"),
}

# Standard PowerPoint scratch layouts that should never be offered to the
# arranger as content layouts (case-insensitive exact match on the name).
_SKIP_CONTENT_NAMES = {"blank", "title and content", "title only"}


def _has_type(text_phs, *types) -> bool:
    return any(_type_matches(p, types) for p in text_phs)


def _type_matches(ph: PlaceholderInfo, types) -> bool:
    return ph.type in {t.name if hasattr(t, "name") else t for t in types}


def _name_hint_role(name: str) -> Optional[str]:
    low = (name or "").lower()
    for role, hints in _ROLE_NAME_HINTS.items():
        if any(h in low for h in hints):
            return role
    return None


def _infer_roles(spec: TemplateSpec) -> None:
    """Populate ``spec.roles`` and ``spec.content_layouts``.

    Strategy: a name hint wins if present and the layout plausibly fits the
    role; otherwise fall back to geometry. Each special role keeps only the
    first match; everything else becomes a content layout.
    """
    roles: Dict[str, Optional[str]] = {"cover": None, "toc": None, "divider": None, "last": None}
    content: List[str] = []

    def text_type_set(layout: LayoutSpec):
        return {p.type for p in layout.text_phs}

    for name, layout in spec.layouts.items():
        types = text_type_set(layout)
        n_text = len(layout.text_phs)
        n_pic = len(layout.pic_phs)
        hint = _name_hint_role(name)

        # --- decide a candidate special role ---
        candidate = None
        if hint == "cover" or ({"CENTER_TITLE", "SUBTITLE"} <= types and n_pic == 0 and "BODY" not in types):
            candidate = "cover"
        elif hint == "last" or ({"CENTER_TITLE"} <= types and n_text == 1 and n_pic == 0):
            candidate = "last"
        elif hint == "divider" or ("CENTER_TITLE" in types and n_pic == 0 and n_text <= 2 and "BODY" in types):
            candidate = "divider"
        elif hint == "toc" or ("BODY" in types and "CENTER_TITLE" not in types and n_pic == 0 and n_text == 1):
            candidate = "toc"

        if candidate and roles.get(candidate) is None:
            roles[candidate] = name
        elif name.strip().lower() in _SKIP_CONTENT_NAMES:
            continue  # standard PPT scratch layout — not a usable content layout
        else:
            content.append(name)

    spec.roles = roles
    # Sort content layouts for stable prompt output: text-only first, then by (n_pic, n_body).
    spec.content_layouts = sorted(
        content, key=lambda nm: (spec.layouts[nm].n_pic, spec.layouts[nm].n_body, nm)
    )


def scan_template(pptx_path) -> TemplateSpec:
    """Parse a ``.pptx`` into a :class:`TemplateSpec`.

    Applies an optional ``<stem>.slidegen.json`` sidecar override.
    """
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))

    spec = TemplateSpec(
        path=str(pptx_path),
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
    )

    for layout in prs.slide_layouts:
        text_phs, pic_phs = _collect_placeholders(layout)
        ls = LayoutSpec(name=layout.name, text_phs=text_phs, pic_phs=pic_phs)
        _assign_content_roles(ls, prs.slide_width)
        spec.layouts[layout.name] = ls

    _infer_roles(spec)
    _apply_sidecar(spec, pptx_path)
    return spec


# --------------------------------------------------------------------------- #
# Sidecar override
# --------------------------------------------------------------------------- #
def _apply_sidecar(spec: TemplateSpec, pptx_path: Path) -> None:
    """Merge an optional ``<stem>.slidegen.json`` override file.

    Schema (all keys optional)::

        {
          "roles": {"cover": "<layout>", "toc": ..., "divider": ..., "last": ...},
          "layouts": {
            "<layout name>": {
              "part_idx": 1, "title_idx": 2,
              "body_idxs": [3,4], "pic_idxs": [5]
            }
          }
        }
    """
    sidecar = pptx_path.with_suffix(".slidegen.json")
    if not sidecar.exists():
        return
    try:
        data = json.loads(sidecar.read_text(encoding="utf-8"))
    except Exception as exc:  # noqa: BLE001 - never let a bad sidecar break scanning
        print(f"[template] WARN: failed to read sidecar {sidecar}: {exc}")
        return

    role_override = data.get("roles") or {}
    for role, layout_name in role_override.items():
        spec.roles[role] = layout_name
        if layout_name in spec.content_layouts:
            spec.content_layouts.remove(layout_name)

    for layout_name, ov in (data.get("layouts") or {}).items():
        ls = spec.layouts.get(layout_name)
        if ls is None:
            continue
        if "part_idx" in ov:
            ls.part_idx = ov["part_idx"]
        if "title_idx" in ov:
            ls.title_idx = ov["title_idx"]
        if "body_idxs" in ov:
            ls.body_idxs = list(ov["body_idxs"])
        if "pic_idxs" in ov:
            ls.pic_idxs = list(ov["pic_idxs"])

    # Recompute content list: special-role layouts must not be selectable.
    role_layouts = {v for v in spec.roles.values() if v}
    spec.content_layouts = [n for n in spec.content_layouts if n not in role_layouts]


# --------------------------------------------------------------------------- #
# Arranger prompt library
# --------------------------------------------------------------------------- #
def _when_to_use(layout: LayoutSpec, slide_width: int) -> str:
    nb, npx = layout.n_body, layout.n_pic
    if npx == 0:
        if nb >= 2:
            return f"Text-only, {nb} text columns/blocks; good for dense or paired text."
        return "Text-only; bullets, no figures."
    parts = [f"{npx} image/figure slot{'s' if npx != 1 else ''}"]
    if nb >= 1:
        parts.append(f"{nb} text block{'s' if nb != 1 else ''}")
    else:
        parts.append("no text block")
    return ", ".join(parts) + "."


def build_layout_library_md(spec: TemplateSpec) -> str:
    """Render the content-layout table injected into the arranger prompt.

    Columns: layout id (the exact ``template_id`` to emit), capacity, guidance.
    """
    lines = [
        "| template_id | capacity | when to use |",
        "|-------------|----------|-------------|",
    ]
    for name in spec.content_layouts:
        ls = spec.layouts[name]
        cap = f"text slots: {ls.n_body}, image slots: {ls.n_pic}"
        lines.append(f"| {name} | {cap} | {_when_to_use(ls, spec.slide_width)} |")

    note = (
        "\nRules:\n"
        "- `template_id` MUST be one of the ids above (exact string).\n"
        "- Count every formula as an image slot: choose a layout whose image "
        "slots >= images + tables + formulas.\n"
        "- If a subsection has no visuals, prefer a text-only layout.\n"
    )
    return "\n".join(lines) + "\n" + note


# --------------------------------------------------------------------------- #
# Template path resolution
# --------------------------------------------------------------------------- #
def resolve_template_path(template_arg, template_dir: str = DEFAULT_TEMPLATE_DIR) -> Path:
    """Locate the template ``.pptx`` to use.

    * ``template_arg`` a stem/filename/path -> resolve against ``template_dir``.
    * legacy int (e.g. ``3``) -> ``slides{int}_template.pptx`` for back-compat.
    * ``None``/empty -> the single ``.pptx`` in the dir, else newest (with a warning).
    """
    tdir = Path(template_dir)

    # Legacy integer template id.
    if isinstance(template_arg, int):
        p = tdir / f"slides{template_arg}_template.pptx"
        if p.exists():
            return p
        raise FileNotFoundError(f"Legacy template not found: {p}")

    if template_arg:
        cand = Path(template_arg)
        tries = [
            cand,                                   # absolute / cwd-relative path
            tdir / cand.name,                       # filename in template dir
            tdir / f"{cand.stem}.pptx",             # stem in template dir
            tdir / f"{template_arg}.pptx",          # raw arg + .pptx
        ]
        for t in tries:
            if t.exists():
                return t
        # case-insensitive fallback
        want = f"{cand.stem}.pptx".lower()
        for f in tdir.glob("*.pptx"):
            if f.name.lower() == want:
                return f
        raise FileNotFoundError(
            f"Template '{template_arg}' not found. Tried: "
            + ", ".join(str(t) for t in tries)
            + f". Available: {[f.name for f in tdir.glob('*.pptx')]}"
        )

    # No arg: pick from the directory.
    pptx_files = sorted(tdir.glob("*.pptx"))
    if not pptx_files:
        raise FileNotFoundError(f"No .pptx template found in {tdir}")
    if len(pptx_files) == 1:
        return pptx_files[0]
    newest = max(pptx_files, key=lambda p: p.stat().st_mtime)
    print(
        f"[template] WARN: multiple templates in {tdir}; using newest '{newest.name}'. "
        f"Pass --template=<name> to choose explicitly."
    )
    return newest


# --------------------------------------------------------------------------- #
# Debug / sanity entry point
# --------------------------------------------------------------------------- #
if __name__ == "__main__":
    import argparse

    p = argparse.ArgumentParser(description="Inspect a SlideGen template .pptx")
    p.add_argument("--template", default="slides3_template")
    p.add_argument("--template_dir", default=DEFAULT_TEMPLATE_DIR)
    a = p.parse_args()

    path = resolve_template_path(a.template, a.template_dir)
    spec = scan_template(path)
    print(f"Template: {spec.path}")
    print(f"Slide size (EMU): {spec.slide_width} x {spec.slide_height}")
    print(f"Roles: {spec.roles}")
    print(f"Content layouts ({len(spec.content_layouts)}):")
    for name in spec.content_layouts:
        ls = spec.layouts[name]
        print(
            f"  - {name}: part={ls.part_idx} title={ls.title_idx} "
            f"body={ls.body_idxs} pic={ls.pic_idxs}"
        )
    print("\n--- arranger library ---")
    print(build_layout_library_md(spec))
