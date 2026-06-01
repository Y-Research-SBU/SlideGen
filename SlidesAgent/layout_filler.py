import json
from pathlib import Path
from typing import Dict, List
import os
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.dml.color import MSO_THEME_COLOR
      
from pprint import pprint

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE   
from pptx.util import Inches

import json
import difflib
from pptx import Presentation
from SlidesAgent.apply_color import *
from SlidesAgent.template_introspect import (
    scan_template,
    resolve_template_path,
    LayoutSpec,
)



COLOR_WHITE = RGBColor(0, 0, 0) 
THEME_COLOR = RGBColor(185, 210, 153) 
 

def _insert_picture_keep_ratio(ph, img_path: Path):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image

    slide = ph.part.slide   
    ph_left, ph_top = ph.left, ph.top
    ph_width, ph_height = ph.width, ph.height
 
    with Image.open(img_path) as img:
        iw, ih = img.size
    aspect = iw / ih
 
    if ph_width / ph_height > aspect:
        new_h = ph_height
        new_w = int(new_h * aspect)
    else:
        new_w = ph_width
        new_h = int(new_w / aspect)
 
    new_left = ph_left + int((ph_width - new_w) / 2)
    new_top = ph_top + int((ph_height - new_h) / 2)
 
    pic = slide.shapes.add_picture(str(img_path), new_left, new_top, width=new_w, height=new_h)
 
    # --- move picture *behind* all text placeholders ---
    spTree = slide.shapes._spTree
    spTree.remove(pic.element)         # temporarily take it out
    spTree.insert(2, pic.element)      # index 0=background,1=layout 

    # --- finally, remove the now‑unused picture placeholder itself ---
    ph.element.getparent().remove(ph.element)
  
from pptx.util import Pt
from PIL import Image

def insert_image_below_content(slide, img_path: Path):
    """
    Insert an image below the lowest existing shape (text or image) on the slide.
    Centered horizontally. Resizes if not enough space.
    """
     
    with Image.open(img_path) as img:
        width_px, height_px = img.size
    aspect_ratio = width_px / height_px

    # slide size（EMUs）
    slide_width = slide.part.slide_layout.part.package.presentation_part.slide_width
    slide_height = slide.part.slide_layout.part.package.presentation_part.slide_height

    
    target_width = slide_width * 0.6
    target_height = target_width / aspect_ratio

     
    lowest_bottom = 0
    for shape in slide.shapes:
        bottom = shape.top + shape.height
        if bottom > lowest_bottom:
            lowest_bottom = bottom

    margin = Pt(20)
    available_space = slide_height - lowest_bottom - margin

    if available_space < target_height:
        target_height = available_space
        target_width = target_height * aspect_ratio
        if target_height <= 0:
            print("Not enough space to insert image:", img_path)
            return

    
    left = (slide_width - target_width) // 2
    top = lowest_bottom + margin

    slide.shapes.add_picture(str(img_path), left, top, width=int(target_width), height=int(target_height))

 

TEXT_TYPES = {
    PH_TYPE.TITLE,
    PH_TYPE.CENTER_TITLE,
    PH_TYPE.SUBTITLE,
    PH_TYPE.BODY,
}

def find_text_placeholders(slide): 
    """Return (part_num_ph, subsection_ph, body_ph) by position."""
    txt_ph = [
        s for s in slide.shapes
        if (
            s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and
            s.has_text_frame and
            s.placeholder_format.type in TEXT_TYPES   
        )
    ]

    if len(txt_ph) < 3:
        raise ValueError("Not enough text placeholders on this slide")

     
    txt_ph.sort(key=lambda s: s.top)
 
    first_row = sorted(txt_ph[:2], key=lambda s: s.left)
    part_ph, title_ph = first_row
    body_ph = txt_ph[2]          

    return part_ph, title_ph, body_ph

def get_content(sec_title, sub_title,outline):
    _, sub = _best_match(outline, sec_title, sub_title)
    return sub.get("content", "") if sub else ""

def set_font_color(paragraph, theme_color):
    paragraph.font.fill.solid()
    if theme_color is not None:
        paragraph.font.fill.fore_color.theme_color = theme_color
    else:
        paragraph.font.fill.fore_color.rgb = RGBColor(255, 105, 180)  # pink fallback

def insert_visuals_auto(slide, visuals: list[Path]):
    """
    Automatically insert visuals (images, tables, formulas) into all available
    picture placeholders on the given slide.
    """
    # Find all picture placeholder shapes
    picture_placeholders = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and
           "Picture" in shape.name
    ]

    if len(visuals) > len(picture_placeholders):
        print(f"Warning: not enough picture placeholders on slide (needed {len(visuals)}, found {len(picture_placeholders)})")
        # remaining = visuals[len(picture_placeholders):]
        # for img_path in remaining:
        #     insert_image_below_content(slide, Path(img_path))

    #  Insert images one by one
    for img_path, ph in zip(visuals, picture_placeholders):
        _insert_picture_keep_ratio(ph, Path(img_path))


def _placeholder_by_name(slide, name: str):
    """Return placeholder shape whose .name == name."""

    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f'Placeholder "{name}" not found on slide master.')
 
 
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def scan_layout_placeholders(template_path: str):
    prs = Presentation(template_path)
    layout_map = {}

    for layout in prs.slide_layouts:
        layout_name = layout.name
        placeholder_names = []

        for shape in layout.shapes:
            shape_type = shape.shape_type
            shape_type_name = str(shape_type)
            try:
                shape_type_name = MSO_SHAPE_TYPE(shape_type).name
            except ValueError:
                pass

            placeholder_names.append({
                "name": shape.name,
                "type": shape_type_name
            })

        layout_map[layout_name] = placeholder_names

    return layout_map


from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_theme_color_from_title(prs, layout_index=2):
    slide_layout = prs.slide_layouts[layout_index]
    temp_slide = prs.slides.add_slide(slide_layout)
 
    for shape in temp_slide.shapes:
        if (
            shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and 
            shape.placeholder_format.type == 1  # 1 == TITLE
        ):
            para = shape.text_frame.paragraphs[0]
            font_color = para.font.color
            if font_color.type == 2:  # 2 == THEME
                theme_color = font_color.theme_color 
                xml_slides = prs.slides._sldIdLst  
                xml_slides.remove(xml_slides[-1])  # remove last slide
                return theme_color
    return None
 
def resolve_visual_paths(slide_info, args):
    import os
    import re
    import json
    from pathlib import Path

    paper = args.paper_name
    prefix = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables"
    base_dir = Path(prefix) / paper
 
    images_json_path = Path(prefix) / f"{paper}_images.json"
    tables_json_path = Path(prefix) / f"{paper}_tables.json"
    images_json = json.load(open(images_json_path, "r", encoding="utf-8")) if images_json_path.exists() else {}
    tables_json = json.load(open(tables_json_path, "r", encoding="utf-8")) if tables_json_path.exists() else {}
 
    if args.formula_mode == 3:
        formulas_dir = Path("contents") / paper / "formula_images"
    else:
        formulas_dir = base_dir
 
    def _norm_digits(s: str) -> str:
 
        m = re.search(r'(\d+)(?!.*\d)', str(s))
        if not m:
            return str(s)
        v = m.group(1).lstrip("0")
        return v if v != "" else "0"

    def _stem_name(p: str) -> str:
        return Path(str(p)).name

    def _extract_figure_numbers(caption: str):
         
        out = set()
        if not caption:
            return out
        caps = str(caption)

        #  Figure 6–7 / Fig. 10-12 / 图 3-5
        for m in re.finditer(r'(?:(?:fig(?:ure)?|图)\.?\s*)(\d+)\s*[–-]\s*(\d+)', caps, flags=re.I):
            a, b = int(m.group(1)), int(m.group(2))
            if a <= b:
                out.update(range(a, b + 1))
            else:
                out.update(range(b, a + 1))

        #  Figure 6 / Fig. 6 / 图6
        for m in re.finditer(r'(?:(?:fig(?:ure)?|图)\.?\s*)(\d+)', caps, flags=re.I):
            out.add(int(m.group(1)))

        return out

    def _extract_table_numbers(caption: str):
         
        out = set()
        if not caption:
            return out
        caps = str(caption)

        for m in re.finditer(r'(?:(?:tab(?:le)?|表)\.?\s*)(\d+)\s*[–-]\s*(\d+)', caps, flags=re.I):
            a, b = int(m.group(1)), int(m.group(2))
            if a <= b:
                out.update(range(a, b + 1))
            else:
                out.update(range(b, a + 1))

        for m in re.finditer(r'(?:(?:tab(?:le)?|表)\.?\s*)(\d+)', caps, flags=re.I):
            out.add(int(m.group(1)))

        return out

    def _build_caption_index_images(images_dict: dict):
         
        idx = {}
        for k, meta in (images_dict or {}).items():
            caps = (meta or {}).get("caption", "")
            for f in _extract_figure_numbers(caps):
                idx.setdefault(f, set()).add(str(k))
        return idx

    def _build_caption_index_tables(tables_dict: dict):
         
        idx = {}
        for k, meta in (tables_dict or {}).items():
            caps = (meta or {}).get("caption", "")
            for t in _extract_table_numbers(caps):
                idx.setdefault(t, set()).add(str(k))
        return idx

    fignum_to_imgkeys = _build_caption_index_images(images_json)
    tbnum_to_tbkeys  = _build_caption_index_tables(tables_json)

     
    images_json = {str(k): v for k, v in images_json.items()}
    tables_json = {str(k): v for k, v in tables_json.items()}

    def _resolve_and_check(path1: Path) -> Path:
        
        if path1.exists():
            return path1
 
        p2 = base_dir / path1.name
        if p2.exists():
            return p2
 
        candidates = [p for p in base_dir.glob("*") if p.name.lower() == path1.name.lower()]
        if candidates:
            return candidates[0]
 
        existing = "\n".join(sorted(p.name for p in base_dir.glob("*.*"))[:80])
        raise FileNotFoundError(
            f"Missing visual file: {path1}\n"
            f"Tried: {path1}, {p2}. Base dir: {base_dir}\n"
            f"Existing (first 80):\n{existing}"
        )

    def _match_by_id_or_caption(img_id: str, name: str, mapping: dict, cap_index: dict, kind: str):
          
        img_id_norm = _norm_digits(img_id)
        name_norm   = _norm_digits(_stem_name(name))
 
        if img_id in mapping:
            return mapping[img_id], img_id
        if img_id_norm in mapping:
            return mapping[img_id_norm], img_id_norm
        if name_norm in mapping:
            return mapping[name_norm], name_norm
 
        if img_id_norm.isdigit():
            num = int(img_id_norm)
            cand = sorted(cap_index.get(num, []), key=lambda x: int(_norm_digits(x)) if _norm_digits(x).isdigit() else 10**9)
            if cand:
                k = cand[0]
                return mapping[k], k
 
        if name_norm.isdigit():
            num2 = int(name_norm)
            cand2 = sorted(cap_index.get(num2, []), key=lambda x: int(_norm_digits(x)) if _norm_digits(x).isdigit() else 10**9)
            if cand2:
                k2 = cand2[0]
                return mapping[k2], k2
 
        avail = sorted(mapping.keys(), key=lambda x: int(_norm_digits(x)) if _norm_digits(x).isdigit() else 10**9)
        raise KeyError(
            f"{kind.capitalize()} id not found. "
            f"given_id='{img_id}', file='{name}', "
            f"norm(id)={img_id_norm}, norm(file)={name_norm}. "
            f"Available keys (first 40): {avail[:40]}"
        )

    # ----------------- images -----------------
    image_paths = []
    for img_str in slide_info.get("images", []):
        name = _stem_name(img_str)
 
        #   'picture|image|fig|table|formula-<num>.ext'
        m = re.search(r'(?:picture|image|fig|table|formula)[-_](\d+)(?=\.[A-Za-z0-9]+$)', name, re.I)
        img_id = m.group(1) if m else _norm_digits(name)

       
        rec, used_key = _match_by_id_or_caption(img_id, name, images_json, fignum_to_imgkeys, kind="image")

        target = rec.get("image_path") or rec.get("path") or rec.get("file")
        if not target:
            raise KeyError(f"Image record has no path field. key={used_key}, record={rec}")

        img_path = _resolve_and_check(Path(target))
        image_paths.append(img_path)

    # ----------------- tables -----------------
    table_paths = []
    for tb_str in slide_info.get("tables", []):
        name = _stem_name(tb_str)
        m = re.search(r'(?:table|tab|picture|image|fig|formula)[-_](\d+)(?=\.[A-Za-z0-9]+$)', name, re.I)
        tb_id = m.group(1) if m else _norm_digits(name)

        rec, used_key = _match_by_id_or_caption(tb_id, name, tables_json, tbnum_to_tbkeys, kind="table")

        target = rec.get("table_path") or rec.get("image_path") or rec.get("path") or rec.get("file")
        if not target:
            raise KeyError(f"Table record has no path field. key={used_key}, record={rec}")

        tb_path = _resolve_and_check(Path(target))
        table_paths.append(tb_path)

    # ----------------- formulas -----------------
    formula_paths = []
    for fname in slide_info.get("formulas", []):
        if args.formula_mode == 3:
            final_path = formulas_dir / fname
        else: 
            final_path = Path(resolve_formula_mode1_path(fname, args))
        formula_paths.append(_resolve_and_check(Path(final_path)))
    print("formula_paths",formula_paths)
    return image_paths + table_paths + formula_paths

 
 


  
import re
from pathlib import Path
import difflib

 
def _extract_idx(val):
   
    if isinstance(val, int):
        return val
    m = re.findall(r"\d+", str(val))
    return int(m[-1]) if m else None

def _nums_from_files(files):
    """['image_7.png', 'table_1.png'] -> {7, 1}"""
    out = set()
    for f in files:
        m = re.findall(r"\d+", str(f))
        if m:
            out.add(int(m[-1]))
    return out

def _best_match(data, sec_title: str, sub_title: str, min_ratio: float = 0.55):
    best_sec, best_sec_score = None, 0.0
    for sec in data.get("sections", []):
        s = difflib.SequenceMatcher(None, sec.get("title","").lower(), (sec_title or "").lower()).ratio()
        if s > best_sec_score:
            best_sec, best_sec_score = sec, s
    if not best_sec or best_sec_score < min_ratio:
        return None, None

    best_sub, best_sub_score = None, 0.0
    for sub in best_sec.get("subsections", []):
        s = difflib.SequenceMatcher(None, sub.get("title","").lower(), (sub_title or "").lower()).ratio()
        if s > best_sub_score:
            best_sub, best_sub_score = sub, s
    if not best_sub or best_sub_score < min_ratio:
        return None, None
    return best_sec, best_sub
 
def _collect_reasons_for_kind(sec, sub, files, figs_data, kind: str):
    
    _, sub_d = _best_match(figs_data, sec, sub)
    if not sub_d:
        return ""

    want = _nums_from_files(files)
    if not want:
        return ""

    #   N -> value_idx
    pairs = []  # [(N, value_idx)]
    for k, v in sub_d.items():
        k_low = str(k).lower()
        if not k_low.startswith(kind):  # 'image' or 'table'
            continue
        N = _extract_idx(k) 
        v_idx = _extract_idx(v)   
        if N is None or v_idx is None:
            continue
        if v_idx in want:
            pairs.append((N, v_idx))

    reasons = []
     
    only_one_asset = (len(pairs) == 1)

    for N, _v_idx in pairs:
         
        candidate_keys = [
            f"reason{N}",
        ]
        if kind == "table":
            candidate_keys.append(f"reasonT{N}")
        else:
            candidate_keys.extend([f"reasonI{N}", f"reasonImg{N}"])
         
        if only_one_asset:
            candidate_keys.append("reason")

         
        found = None
        for rk in candidate_keys:
            if rk in sub_d and isinstance(sub_d[rk], str) and sub_d[rk].strip():
                found = sub_d[rk].strip()
                break

         
        if not found and "reason" in sub_d and isinstance(sub_d["reason"], str) and sub_d["reason"].strip():
            found = sub_d["reason"].strip()

        if found:
            reasons.append(found)

    return "\n".join(reasons)

def get_image_reasons(sec, sub, image_files, figs_data):
    return _collect_reasons_for_kind(sec, sub, image_files, figs_data, kind="image")

def get_table_reasons(sec, sub, table_files, figs_data):
    return _collect_reasons_for_kind(sec, sub, table_files, figs_data, kind="table")

def get_formula_reasons(sec, sub, formula_files, formula_data) -> str:
     
    _, sub_d = _best_match(formula_data, sec, sub)
    if not sub_d:
        return ""
 
    pairs = []
    for k, v in sub_d.items():
        k_low = str(k).lower()
        if not k_low.startswith("formula"):
            continue
        N = _extract_idx(k)   
        if N is None:
            continue
        rkey = f"reason{N}"
        rtxt = sub_d.get(rkey)
        if isinstance(rtxt, str) and rtxt.strip():
            pairs.append((N, rtxt.strip()))

    
    try:
        pairs.sort(key=lambda x: int(x[0]))
    except Exception: 
        pass

   
    if not pairs and isinstance(sub_d.get("reason"), str) and sub_d["reason"].strip():
        if formula_files and len(formula_files) == 1:
            return sub_d["reason"].strip()
 
    if formula_files and len(pairs) > len(formula_files):
        pairs = pairs[:len(formula_files)]
 
    out, seen = [], set()
    for _, r in pairs:
        if r not in seen:
            out.append(r)
            seen.add(r)

    return "\n".join(out)

def resolve_formula_mode1_path(fname: str, args) -> Path:
    """
    Extract formula index i from fname (like "formula_4.png"),
    and generate the path:
    <{args.model_name_t}_{args.model_name_v}>_images_and_tables/{args.paper_name}/{args.paper_name}-formula-i.png
    """ 

    stem = Path(fname).stem   
    match = re.search(r'(\d+)(?!.*\d)', stem) 
    # match = re.search(r'(?i)\bformula(?:[_\-\s]*)?(\d+)\b', stem)
    if not match:
        raise ValueError(f"Cannot extract index from formula filename: {fname}")
    
    i = match.group(1)
    path_str = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables/{args.paper_name}/{args.paper_name}-formula-{i}.png"
    return Path(path_str)

import json
from pathlib import Path
from typing import Dict, Any, List

def _is_T1_textonly(s: Dict[str, Any]) -> bool:
     
    return (
        s.get("template_id") == "T1_TextOnly"
        and not s.get("images")
        and not s.get("tables")
        and not s.get("formulas")
    )

def pair_T1_to_T19(plan_path: str, write_back: bool = True) -> int:
 
    p = Path(plan_path)
    plan = json.loads(p.read_text(encoding="utf-8"))
    slides: List[Dict[str, Any]] = plan.get("slides", [])
    out: List[Dict[str, Any]] = []

    i, n, made = 0, len(slides), 0
    while i < n:
        cur = slides[i]
        if (
            i + 1 < n
            and _is_T1_textonly(cur)
            and _is_T1_textonly(slides[i + 1])
            and cur.get("section") == slides[i + 1].get("section")
        ):
            left, right = cur, slides[i + 1]
            out.append({
                "section": cur.get("section"),
                "template_id": "T19_2Text",
                "columns": [
                    {
                        "subsection": left.get("subsection", "") or "",
                        "bullets": left.get("bullets", []) or []
                    },
                    {
                        "subsection": right.get("subsection", "") or "",
                        "bullets": right.get("bullets", []) or []
                    }
                ]
            })
            made += 1
            i += 2
            continue
 
        out.append(cur)
        i += 1

    plan["slides"] = out
    if write_back:
        p.write_text(json.dumps(plan, ensure_ascii=False, indent=4), encoding="utf-8")
    return made

def validate_no_consecutive_T1(plan_path: str) -> List[int]:
     
    p = Path(plan_path)
    plan = json.loads(p.read_text(encoding="utf-8"))
    slides: List[Dict[str, Any]] = plan.get("slides", [])
    bad_idxs: List[int] = []
    for i in range(len(slides) - 1):
        a, b = slides[i], slides[i + 1]
        if _is_T1_textonly(a) and _is_T1_textonly(b) and a.get("section") == b.get("section"):
            bad_idxs.append(i)
    return bad_idxs
from pptx.util import Pt

def _clear_text_frame(tf):
    if tf.paragraphs:
        tf.paragraphs[0].text = ""
         
        while len(tf.paragraphs) > 1:
            p = tf._element.p_lst[-1]
            p.getparent().remove(p)
    else:
        tf.clear()

def _fill_bullets(tf, bullets, lvl0_size=24, lvl1_size=24):
    _clear_text_frame(tf)
    first = True
    for b in (bullets or []):
        # Reuse the (empty) first paragraph instead of leaving a blank line on top.
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = (b.get("text") or "").strip()
        p.level = 0
        p.font.size = Pt(lvl0_size)

        for s in (b.get("sub") or []):
            sp = tf.add_paragraph()
            sp.text = str(s).strip()
            sp.level = 1
            sp.font.size = Pt(lvl1_size)
from pptx.enum.shapes import PP_PLACEHOLDER

def _get_placeholder(slide, name): 
    for shp in slide.shapes:
        if getattr(shp, "name", "").strip().lower() == name.strip().lower():
            return shp
    return None

def _ph_text_n(slide, n:int): 
    targets = {f"text placeholder {n}", f"文本占位符 {n}"}
    for shp in slide.shapes:
        nm = getattr(shp, "name", "").strip().lower()
        if nm in targets:
            return shp 
        if ("placeholder" in nm or "占位符" in nm) and nm.endswith(f" {n}"):
            return shp
    return None
def _ph_by_idx(slide, idx:int):
    for ph in getattr(slide, "placeholders", []):
        if ph.placeholder_format.idx == idx:
            return ph
    return None

def fill_T19_2Text(slide, slide_info, section_no_text):
     
    print("slide_info:")
    print(slide_info)
    part_ph = (
        _get_placeholder(slide, "Part")
        or _get_placeholder(slide, "Text Placeholder 2")
        or _ph_text_n(slide, 2)
    )
   
    title_bar = _get_placeholder(slide, "Text Placeholder 1") or _ph_text_n(slide, 1)
    title_bar = _ph_by_idx(slide, 1)  or _ph_text_n(slide, 1)

     
    section_title = slide_info.get("section") 
    
    title_bar.text = section_title
    tf = title_bar.text_frame
    tf.clear()                        
    tf.paragraphs[0].text = section_title
    print("[AFTER] title_bar text =", repr(tf.text))

    lt = (
        _get_placeholder(slide, "Left Title")
        or _get_placeholder(slide, "Text Placeholder 3")
        or _ph_text_n(slide, 3)
    )
    lb = (
        _get_placeholder(slide, "Left Body")
        or _get_placeholder(slide, "Text Placeholder 4")
        or _ph_text_n(slide, 4)
    )
    rt = (
        _get_placeholder(slide, "Right Title")
        or _get_placeholder(slide, "Text Placeholder 5")
        or _ph_text_n(slide, 5)
    )
    rb = (
        _get_placeholder(slide, "Right Body")
        or _get_placeholder(slide, "Text Placeholder 6")
        or _ph_text_n(slide, 6)
    )

    
    cols = slide_info.get("columns") or []
    left  = cols[0] if len(cols) > 0 else {}
    right = cols[1] if len(cols) > 1 else {}
 
    if part_ph is not None and getattr(part_ph, "has_text_frame", False):
        part_ph.text_frame.text = f"{section_no_text}"

    if title_bar is not None and getattr(title_bar, "has_text_frame", False):
        title_txt = slide_info.get("section", "") or slide_info.get("title", "")
        title_bar.text_frame.text = title_txt

    if lt is not None and getattr(lt, "has_text_frame", False):
        lt.text_frame.text = left.get("subsection", "") or left.get("title", "") or ""

    if rt is not None and getattr(rt, "has_text_frame", False):
        rt.text_frame.text = right.get("subsection", "") or right.get("title", "") or ""

    if lb is not None and getattr(lb, "has_text_frame", False):
        _fill_bullets(lb.text_frame, left.get("bullets"))

    if rb is not None and getattr(rb, "has_text_frame", False):
        _fill_bullets(rb.text_frame, right.get("bullets"))

     
    missing = []
    for name, ph in [
        ("Part(2)", part_ph), ("Left Title(3)", lt), ("Left Body(4)", lb),
        ("Right Title(5)", rt), ("Right Body(6)", rb),
    ]:
        if ph is None:
            missing.append(name)
    if missing:
        print(f"[WARN] T19_2Text 模板缺少占位：{', '.join(missing)}")



def delete_slide(prs: Presentation, slide_index: int) -> None:
    """
    Delete slide by index in python-pptx (works by removing sldId and dropping rel).
    """
    sldIdLst = prs.slides._sldIdLst  # pylint: disable=protected-access
    sldId_elems = list(sldIdLst)

    if slide_index < 0 or slide_index >= len(sldId_elems):
        raise IndexError(f"slide_index out of range: {slide_index}, total={len(sldId_elems)}")

    sldId = sldId_elems[slide_index]
    rId = sldId.rId
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)


# debug_list_placeholders(slide)

# --------------------------------------------------------------------------- #
# Generic, template-agnostic filling helpers (SlideGen v2)
# --------------------------------------------------------------------------- #
# Default font sizes (pt) used when a slide carries no per-slide "font" block.
DEFAULT_FONTS = {"title": None, "body": 24, "sub": 24, "toc": 36}


def _ph_by_idx_strict(slide, idx):
    """Return the placeholder shape with the given idx, or None."""
    if idx is None:
        return None
    for ph in getattr(slide, "placeholders", []):
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_ph_text(slide, idx, text, font_size=None):
    """Set a single line of text on the placeholder identified by idx."""
    ph = _ph_by_idx_strict(slide, idx)
    if ph is None or not getattr(ph, "has_text_frame", False):
        return False
    tf = ph.text_frame
    tf.clear()
    tf.paragraphs[0].text = str(text or "")
    if font_size:
        tf.paragraphs[0].font.size = Pt(font_size)
    return True


def _fill_bullets_idx(slide, idx, bullets, body_size=24, sub_size=24):
    """Fill a body placeholder (by idx) with hierarchical bullets."""
    ph = _ph_by_idx_strict(slide, idx)
    if ph is None or not getattr(ph, "has_text_frame", False):
        return False
    _fill_bullets(ph.text_frame, bullets, lvl0_size=body_size, lvl1_size=sub_size)
    return True


def _bullets_from_slide_info(slide_info):
    """Normalize a slide_info into a list of column dicts {subsection, bullets}.

    Single-body slides yield one column; multi-column slides (T19-style) use the
    explicit "columns" key when present, otherwise the flat "bullets" list.
    """
    cols = slide_info.get("columns")
    if cols:
        return [
            {
                "subsection": c.get("subsection", "") or c.get("title", "") or "",
                "bullets": c.get("bullets", []) or [],
            }
            for c in cols
        ]
    return [
        {
            "subsection": slide_info.get("subsection", "") or "",
            "bullets": slide_info.get("bullets", []) or [],
        }
    ]


def _distribute_bullets(columns, n_body):
    """Map content columns onto the layout's body placeholders.

    Returns a list of length ``n_body`` of {subsection, bullets} dicts. When the
    layout has more body slots than content columns, extra slots get blank
    content; when fewer, surplus columns are merged into the last slot.
    """
    if n_body <= 0:
        return []
    out = [{"subsection": "", "bullets": []} for _ in range(n_body)]
    if not columns:
        return out
    if len(columns) <= n_body:
        for i, c in enumerate(columns):
            out[i] = c
    else:
        for i in range(n_body - 1):
            out[i] = columns[i]
        merged = {"subsection": columns[n_body - 1].get("subsection", ""), "bullets": []}
        for c in columns[n_body - 1:]:
            merged["bullets"].extend(c.get("bullets", []))
        out[n_body - 1] = merged
    return out


def fill_content_slide(slide, layout_spec: LayoutSpec, slide_info, section_no, fonts=None):
    """Fill a content slide generically from its inferred layout roles.

    * ``part``  -> section number (e.g. "03")
    * ``title`` -> subsection (or section) title
    * ``body_idxs`` -> bullet columns (single or multi-column)

    Honors per-slide font overrides via the ``fonts`` dict
    ({"title","body","sub"}); falls back to DEFAULT_FONTS.
    """
    fonts = {**DEFAULT_FONTS, **(slide_info.get("font") or {}), **(fonts or {})}

    # Part number (tiny top-left box).
    if layout_spec.part_idx is not None:
        _set_ph_text(slide, layout_spec.part_idx, f"{section_no}")

    # Title bar.
    title_txt = slide_info.get("subsection") or slide_info.get("section") or slide_info.get("title") or ""
    if layout_spec.title_idx is not None:
        _set_ph_text(slide, layout_spec.title_idx, title_txt, font_size=fonts.get("title"))

    # Body columns. Each layout column is {"title": idx|None, "body": idx}; a
    # short heading box (when present) gets the column subsection, and bullets go
    # into the tall box — so two-column layouts (T19) keep headings separate
    # from their bullet text instead of dumping everything in one box.
    body_cols = layout_spec.body_cols or [{"title": None, "body": i} for i in layout_spec.body_idxs]
    columns = _bullets_from_slide_info(slide_info)
    mapped = _distribute_bullets(columns, len(body_cols))
    multi = len(body_cols) > 1

    for slot, col in zip(body_cols, mapped):
        bullets = col.get("bullets", [])
        subsec = col.get("subsection", "")
        if slot.get("title") is not None:
            # dedicated heading box for this column
            if multi and subsec:
                _set_ph_text(slide, slot["title"], subsec, font_size=fonts.get("title"))
            _fill_bullets_idx(slide, slot["body"], bullets,
                              body_size=fonts.get("body", 24), sub_size=fonts.get("sub", 24))
        else:
            # no heading box: prefix the subsection as a bullet for multi-col layouts
            if multi and subsec:
                bullets = [{"text": subsec, "sub": []}] + bullets
            _fill_bullets_idx(slide, slot["body"], bullets,
                              body_size=fonts.get("body", 24), sub_size=fonts.get("sub", 24))


def generate_pptx_from_plan(
    args,
    template: Path | int = None
):


    figs_json_path  =  f"contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_figures.json"
    formula_json_path = f"contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_formula_match.json"
    paper_outline_json = f'contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_raw_content.json' 
    with open(paper_outline_json, "r", encoding="utf-8") as f: outline_json  = json.load(f)
    with open(figs_json_path, encoding="utf-8") as f: figs_data   = json.load(f)
    with open(formula_json_path, encoding="utf-8") as f: formula_data   = json.load(f)
    
    plan_json = f'contents/{args.paper_name}/<{args.model_name_t}_{args.model_name_v}>_slide_plan.json'

    # ---------- resolve & scan the template ----------
    if template is None:
        template_path = getattr(args, "template_path", None)
        if not template_path:
            template_path = resolve_template_path(
                getattr(args, "template_name", None),
                getattr(args, "template_dir", "utils/slides_template"),
            )
    else:
        template_path = resolve_template_path(template)  # int or name -> path
    template_path = str(template_path)
    spec = scan_template(template_path)
    print(f"[filler] template={template_path}")
    print(f"[filler] roles={spec.roles}")

    # T1->T19 pairing only makes sense when both layouts exist in this deck.
    if "T1_TextOnly" in spec.layouts and "T19_2Text" in spec.layouts:
        made = pair_T1_to_T19(plan_json)
        print(f"[plan] T1->T19 pairs made: {made}")

    plan: Dict = json.loads(Path(plan_json).read_text(encoding="utf-8"))

    title    = outline_json["metadata"]["title"]
    subtitle = outline_json["metadata"]["author"]

    prs = Presentation(template_path)

    theme_color = extract_theme_color_from_title(prs)
    print("Theme color:", theme_color)

    print("Available slide layouts in template:")
    for layout in prs.slide_layouts:
        print("-", layout.name)

    # ---------- cover ----------
    cover_name = spec.role_layout("cover")
    if cover_name:
        cover_spec = spec.get(cover_name)
        cover = prs.slides.add_slide(prs.slide_layouts.get_by_name(cover_name))
        c_title = cover_spec.title_type_idx()
        if c_title is not None:
            _set_ph_text(cover, c_title, title)
        sub_idxs = cover_spec.secondary_text_idxs(exclude=c_title)
        if sub_idxs:
            _set_ph_text(cover, sub_idxs[0], subtitle)
    else:
        print("[filler] WARN: no cover layout found; skipping cover slide.")

    # ---------- Contents (table of contents) ----------
    toc_name = spec.role_layout("toc")
    seen = set()
    unique_sections = []
    for slide in plan["slides"]:
        sec = slide["section"]
        if sec not in seen:
            seen.add(sec)
            unique_sections.append(sec)

    if toc_name:
        toc_spec = spec.get(toc_name)
        outline = prs.slides.add_slide(prs.slide_layouts.get_by_name(toc_name))
        toc_idx = toc_spec.body_idxs[0] if toc_spec.body_idxs else toc_spec.title_type_idx()
        toc_ph = _ph_by_idx_strict(outline, toc_idx)
        if toc_ph is not None and toc_ph.has_text_frame:
            tf = toc_ph.text_frame
            tf.clear()
            for i, sec in enumerate(unique_sections):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = sec
                p.level = 0
                p.font.size = Pt(DEFAULT_FONTS["toc"])
    else:
        print("[filler] WARN: no TOC layout found; skipping contents slide.")

    divider_name = spec.role_layout("divider")

    # ---------- Body ----------
    current_section  = None
    section_counter  = 0
    for slide_info in plan["slides"]:
        # ---------- Section divider ----------
        if slide_info["section"] != current_section:
            current_section = slide_info["section"]
            section_counter += 1
            if divider_name:
                div_spec = spec.get(divider_name)
                sec_slide = prs.slides.add_slide(prs.slide_layouts.get_by_name(divider_name))
                d_title = div_spec.title_type_idx()
                if d_title is not None:
                    _set_ph_text(sec_slide, d_title, current_section)
                part_idxs = div_spec.secondary_text_idxs(exclude=d_title)
                if part_idxs:
                    _set_ph_text(sec_slide, part_idxs[0], f"PART {section_counter:02d}")

        template_id = slide_info["template_id"]
        layout = prs.slide_layouts.get_by_name(template_id)
        if layout is None:
            raise ValueError(f" Template layout '{template_id}' not found in template '{template_path}'.")
        layout_spec = spec.get(template_id)
        slide = prs.slides.add_slide(layout)

        # ---------- text (generic, role-based) ----------
        fill_content_slide(slide, layout_spec, slide_info, f"{section_counter:02d}")

        # ---------- visuals ----------
        visuals = resolve_visual_paths(slide_info, args)
        insert_visuals_auto(slide, visuals)

        # ---------- note ----------
        notes_chunks = []
        txt = get_content(slide_info["section"], slide_info.get("subsection", ""), outline_json)
        if txt: notes_chunks.append(txt)
        if slide_info.get("images"):
            img_r = get_image_reasons(slide_info["section"], slide_info.get("subsection", ""),
                                      slide_info["images"], figs_data)
            notes_chunks.append(img_r)
        if slide_info.get("tables"):
            tb_r = get_table_reasons(slide_info["section"], slide_info.get("subsection", ""),
                                     slide_info["tables"], figs_data)
            notes_chunks.append(tb_r)
        if slide_info.get("formulas"):
            fm_r = get_formula_reasons(
                slide_info["section"],
                slide_info.get("subsection", ""),
                slide_info["formulas"],
                formula_data  ,
                )
            notes_chunks.append(fm_r)

        if notes_chunks:
            nframe = slide.notes_slide.notes_text_frame
            if nframe.text and not nframe.text.endswith("\n"):
                nframe.text += "\n"
            nframe.text += "\n\n".join([c for c in notes_chunks if c])

    # ---------- thanks / last page ----------
    last_name = spec.role_layout("last")
    if last_name:
        last_spec = spec.get(last_name)
        thanks = prs.slides.add_slide(prs.slide_layouts.get_by_name(last_name))
        l_title = last_spec.title_type_idx()
        if l_title is not None:
            _set_ph_text(thanks, l_title, "THANKS!")
            tph = _ph_by_idx_strict(thanks, l_title)
            if tph is not None and tph.text_frame.paragraphs[0].runs:
                tph.text_frame.paragraphs[0].runs[0].font.bold = True
    else:
        print("[filler] WARN: no last-page layout found; skipping thanks slide.")

    output_pptx = f'contents/{args.paper_name}/{args.model_name_t}_{args.model_name_v}_output_slides.pptx'
    prs.save(str(output_pptx))
    delete_slide(prs, 0)
    prs.save(str(output_pptx))

    prefix = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables"
    base_dir = Path(prefix) / args.paper_name
    target_name = f"{args.paper_name}-with-image-refs_artifacts"
    
    artifacts_dir = base_dir / target_name
    
    if not artifacts_dir.exists():
        hits = list(base_dir.glob(f"**/{target_name}"))
        if hits: 
            artifacts_dir = max(hits, key=lambda p: len(p.parts))

    img_candidates = sorted(artifacts_dir.glob("image_*.png"))
    theme_imgs = [str(p) for p in img_candidates[:2]]


    if len(theme_imgs) < 1:
        print(f"[warn] No image_*.png found under: {artifacts_dir}. Skip theming.")
    else:
 
        theme_hex, base_hex = pick_theme_color(
            images=theme_imgs,
            prefer_dark=True,
            min_v=0.10,
            max_v=0.99,
            return_base_hex=True,
        )
        print("[theme] base_hex :", base_hex)
        print("[theme] theme_hex:", theme_hex)
        print("[theme] imgs:", theme_imgs)

        
        output_pptx_path = Path(output_pptx)
        themed_pptx_path = output_pptx_path.with_name(output_pptx_path.stem + "_themed" + output_pptx_path.suffix)
        print("colored path")
       
        set_one_theme_color(
            pptx_in=str(output_pptx_path),
            pptx_out=str(themed_pptx_path),
            color_hex=theme_hex,
            target_key="dk2",    
        )

        print(f"[ok] themed pptx saved: {themed_pptx_path}")


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Fill PPTX from slides plan.")
    parser.add_argument(
        "--plan",  
        default="SlideGen/contents/STEP_A_General_and_Scalable_Framework_for_Solving_Video_Inverse_Problems/<4o_4o>_slide_plan.json",  
        help="slides_plan.json"
        )
    parser.add_argument(
        "--paper_name",  
        default="STEP_A_General_and_Scalable_Framework_for_Solving_Video_Inverse_Problems"  
        )
    parser.add_argument(
        "--template",  
        type=int,
        default=3,
        help="Template number, e.g. 3 for slides3_template.pptx"
    )
    parser.add_argument(
        "--out", 
        default="output.pptx" 
        )
    parser.add_argument(
        "--model_name_t", 
        default="4o" 
        )
    parser.add_argument(
        "--model_name_v", 
        default="4o" 
        )
        
    args = parser.parse_args()

    layout_info = scan_layout_placeholders("SlideGen/utils/slides_template/slides3_template.pptx")

    for layout_name, shapes in layout_info.items():
        print(f"\n Layout: {layout_name}")
        for s in shapes:
            print(f"  - {s['name']} ({s['type']})")

 

    generate_pptx_from_plan(args.plan, args.template, args.out)
     
    print(f" Saved to {args.out}")
